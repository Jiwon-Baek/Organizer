#!/usr/bin/env python3
"""
Organizer 앱 사용 설명서 PPTX 생성 스크립트
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

# ── 색상 ─────────────────────────────────────────
RED    = RGBColor(0xCC, 0x22, 0x22)
DARK   = RGBColor(0x1F, 0x1D, 0x1A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BEIGE  = RGBColor(0xF5, 0xF1, 0xE8)
ACCENT = RGBColor(0xB3, 0x54, 0x1E)
MUTED  = RGBColor(0x73, 0x6B, 0x5F)
BLUE   = RGBColor(0x22, 0x55, 0xAA)
GREEN  = RGBColor(0x22, 0x88, 0x22)
ORANGE = RGBColor(0xDD, 0x88, 0x00)
CREAM  = RGBColor(0xE8, 0xD8, 0xC8)

# ── 스크린샷 경로 ─────────────────────────────────
SS_DIR = "/mnt/c/EunYoungKim/Organizer/screenshots/"
SS = {
    "dashboard": SS_DIR + "20260424_163309.png",
    "context":   SS_DIR + "20260424_163324.png",
    "chapter":   SS_DIR + "20260424_163327.png",
    "note":      SS_DIR + "20260424_163330.png",
}

# 실제 해상도 읽기
SS_DIMS = {}
for key, path in SS.items():
    with PILImage.open(path) as im:
        SS_DIMS[key] = im.size  # (w, h)

# ── 슬라이드 설정 ─────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

SS_T = Inches(0.52)   # 슬라이드 상단 제목 아래 스크린샷 시작 위치

# ── 헬퍼 함수 ─────────────────────────────────────

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bg(slide, color=BEIGE):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def add_image(slide, key):
    sw, sh = SS_DIMS[key]
    img_w = W
    img_h = int(W * sh / sw)
    return slide.shapes.add_picture(SS[key], 0, SS_T, img_w, img_h)

def ss_emu(key, px, py, pw, ph):
    """스크린샷 픽셀 좌표 → 슬라이드 EMU 좌표 변환."""
    nw, nh = SS_DIMS[key]
    sw, sh = W, int(W * nh / nw)
    l = int(0   + px / nw * sw)
    t = int(SS_T + py / nh * sh)
    w = int(pw / nw * sw)
    h = int(ph / nh * sh)
    return l, t, w, h

def add_text(slide, text, l, t, w, h, size=13, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, bg_col=None, italic=False):
    box = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    # 여러 줄 처리
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "맑은 고딕"
    if bg_col:
        box.fill.solid()
        box.fill.fore_color.rgb = bg_col
    return box

def add_spot(slide, l, t, w, h, color=RED, lw=Pt(2.5)):
    """투명 배경 + 컬러 테두리 (스포트라이트 박스)."""
    shape = slide.shapes.add_shape(1, int(l), int(t), max(int(w), 1), max(int(h), 1))
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = int(lw)
    return shape

def add_label(slide, text, l, t, w, size=10, bg_col=RED, fg=WHITE, h=None):
    """채워진 레이블 박스 (설명 태그)."""
    if h is None:
        # 줄수 기반 높이 계산
        lines = text.count("\n") + 1
        h = int(Pt(size) * 2.1 * lines)
    shape = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_col
    shape.line.color.rgb = bg_col
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = fg
        r.font.name = "맑은 고딕"
    return shape

def add_card(slide, l, t, w, h, fill=WHITE, border=MUTED):
    """흰색 카드 박스."""
    shape = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = int(Pt(1.0))
    return shape

def add_hline(slide, l, t, w, color=MUTED, lw=Pt(1.0)):
    shape = slide.shapes.add_shape(1, int(l), int(t), int(w), int(Pt(1.0)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════
# SLIDE 1 : 표지
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)

# 상단 장식 스트립
strip = s.shapes.add_shape(1, 0, int(Inches(2.3)), W, int(Inches(2.8)))
strip.fill.solid()
strip.fill.fore_color.rgb = CREAM
strip.line.fill.background()

add_text(s, "Organizer",
         Inches(1), Inches(2.45), Inches(11.33), Inches(1.3),
         size=56, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(s, "사용 설명서",
         Inches(1), Inches(3.65), Inches(11.33), Inches(0.85),
         size=30, color=ACCENT, align=PP_ALIGN.CENTER)

# 구분선
line = s.shapes.add_shape(1, int(Inches(4.5)), int(Inches(4.6)), int(Inches(4.33)), int(Pt(2)))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text(s, "PDF 기반 연구자료 관리 데스크탑 앱",
         Inches(1), Inches(4.8), Inches(11.33), Inches(0.5),
         size=15, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Notebook  →  Chapter  →  Leaf page  3단계 계층 구조",
         Inches(1), Inches(5.35), Inches(11.33), Inches(0.5),
         size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "• 연구자료 분류  • PDF 뷰어 내장  • 페이지별 코멘트  • 전체 검색",
         Inches(1), Inches(5.9), Inches(11.33), Inches(0.5),
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 2 : 전체 구조 개요
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)

add_text(s, "전체 구조 개요",
         Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
         size=22, bold=True, color=DARK)

# 세 박스
bw, bh = Inches(3.6), Inches(4.2)
by = Inches(0.9)
bxs = [Inches(0.5), Inches(4.87), Inches(9.23)]
bcols = [RGBColor(0xE8, 0xD5, 0xC0), RGBColor(0xD0, 0xE8, 0xD0), RGBColor(0xC8, 0xD8, 0xF0)]
btitles = ["📓  Notebook", "📁  Chapter", "📄  Leaf page"]
bdescs = [
    ("최상위 분류 단위\n\n"
     "연구 분야별 대분류\n"
     "예: 상담자료 / 연구자료 / 강의록\n\n"
     "특징\n"
     "• 이모지 아이콘 지정 가능\n"
     "• 이름 변경 가능\n"
     "• 여러 Chapter 포함"),
    ("Notebook 하위 분류\n\n"
     "주제 또는 연도별 중분류\n"
     "예: ADHD / 우울증 / 2026\n\n"
     "특징\n"
     "• 이모지 아이콘 지정 가능\n"
     "• 이름 변경 가능\n"
     "• 여러 Leaf page 포함"),
    ("실제 자료 단위\n\n"
     "개별 주제의 자료 묶음\n"
     "예: 최신논문리뷰 / 증례보고\n\n"
     "포함 정보\n"
     "• 메모 (자유 텍스트)\n"
     "• 태그\n"
     "• PDF 파일 (여러 개 연결)\n"
     "• PDF 페이지별 코멘트"),
]

for i in range(3):
    box = s.shapes.add_shape(1, int(bxs[i]), int(by), int(bw), int(bh))
    box.fill.solid()
    box.fill.fore_color.rgb = bcols[i]
    box.line.color.rgb = MUTED
    box.line.width = int(Pt(1.2))
    add_text(s, btitles[i],
             bxs[i]+Inches(0.18), by+Inches(0.15), bw-Inches(0.36), Inches(0.5),
             size=17, bold=True, color=DARK)
    add_hline(s, bxs[i]+Inches(0.18), by+Inches(0.65), bw-Inches(0.36), color=MUTED)
    add_text(s, bdescs[i],
             bxs[i]+Inches(0.18), by+Inches(0.75), bw-Inches(0.36), Inches(3.3),
             size=12, color=DARK)

# 화살표
for i in range(2):
    ax = bxs[i] + bw + Inches(0.13)
    ay = by + bh/2 - Inches(0.22)
    add_text(s, "→", ax, ay, Inches(0.5), Inches(0.45),
             size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s,
    "💡  사이드바에서 Notebook과 Chapter를 탐색하고, "
    "Chapter를 클릭하면 해당 Leaf page 목록이 메인 화면에 표시됩니다.",
    Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.6),
    size=12, color=MUTED, italic=True)


# ═══════════════════════════════════════════════════════
# SLIDE 3 : 대시보드
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "① 대시보드  —  앱 시작 화면",
         Inches(0.2), Inches(0.04), Inches(8), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "dashboard")

def d(px, py, pw, ph): return ss_emu("dashboard", px, py, pw, ph)

# 사이드바
l,t,w,h = d(0, 0, 214, 768)
add_spot(s, l, t, w, h)
add_label(s, "① 사이드바", l+w+Pt(6), t+Inches(0.1), Inches(1.3))

# 최근 열람 노트
l,t,w,h = d(220, 28, 382, 188)
add_spot(s, l, t, w, h)
add_label(s, "② 최근 열람 노트", l, t-Pt(24), Inches(2.1))

# 최근 열람 PDF
l,t,w,h = d(612, 28, 382, 188)
add_spot(s, l, t, w, h)
add_label(s, "③ 최근 열람 PDF", l, t-Pt(24), Inches(2.1))

# 최근 주요 노트
l,t,w,h = d(220, 224, 382, 188)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "④ 최근 주요 노트", l, t-Pt(24), Inches(2.1), bg_col=BLUE)

# 관련 항목
l,t,w,h = d(612, 224, 382, 188)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "⑤ 관련 항목", l, t-Pt(24), Inches(1.6), bg_col=BLUE)

# 검색창
l,t,w,h = d(1095, 2, 266, 26)
add_spot(s, l, t, w, h, color=GREEN)
add_label(s, "⑥ 검색창 (Enter)", l-Inches(0.3), t+h+Pt(5), Inches(2.0), bg_col=GREEN)

# 언어 선택
l,t,w,h = d(1334, 2, 30, 24)
add_spot(s, l, t, w, h, color=ORANGE)
add_label(s, "KO/EN", l-Inches(0.5), t+h+Pt(5), Inches(1.0), bg_col=ORANGE)


# ═══════════════════════════════════════════════════════
# SLIDE 4 : 사이드바 탐색
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "② 사이드바 탐색  —  Notebook / Chapter 구조",
         Inches(0.2), Inches(0.04), Inches(9), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "chapter")

def c(px, py, pw, ph): return ss_emu("chapter", px, py, pw, ph)

# Notebook 항목
l,t,w,h = c(0, 115, 212, 30)
add_spot(s, l, t, w, h)
add_label(s, "Notebook 항목 (클릭: 접기/펼치기)", l+w+Pt(6), t, Inches(3.1))

# Chapter 항목들
l,t,w,h = c(18, 196, 190, 62)
add_spot(s, l, t, w, h)
add_label(s, "Chapter 항목 (클릭: Leaf page 목록 열기)", l+w+Pt(6), t, Inches(3.4))

# ··· 버튼
l,t,w,h = c(182, 198, 22, 22)
add_spot(s, l, t, w, h, color=ORANGE, lw=Pt(3.0))
add_label(s, "··· 버튼 (마우스 오버시 표시)", l-Inches(2.0), t+h+Pt(6), Inches(2.5), bg_col=ORANGE)

# + 새 Notebook
l,t,w,h = c(5, 843, 204, 26)
add_spot(s, l, t, w, h, color=GREEN)
add_label(s, "+ 새 Notebook 추가", l, t+h+Pt(5), Inches(1.9), bg_col=GREEN)

# 사이드바 접기 버튼
l,t,w,h = c(208, 395, 10, 50)
add_spot(s, l, t, Inches(0.12), h, color=BLUE)
add_label(s, "◀ 사이드바 접기", int(l-Inches(1.6)), t, Inches(1.7), bg_col=BLUE)

# 우측 설명 카드
add_card(s, Inches(8.2), Inches(0.6), Inches(5.0), Inches(6.7))
add_text(s,
    "사이드바 구성\n",
    Inches(8.4), Inches(0.75), Inches(4.7), Inches(0.45),
    size=14, bold=True, color=DARK)
add_text(s,
    "Notebook 항목\n"
    "• 클릭하면 하위 Chapter 목록 접기/펼치기\n"
    "• 마우스 오버 → ··· 버튼 나타남\n\n"
    "Chapter 항목\n"
    "• 클릭하면 메인 화면에 Leaf page 목록 표시\n"
    "• 현재 선택된 항목은 배경색으로 강조됨\n"
    "• 마우스 오버 → ··· 버튼 나타남\n\n"
    "··· 버튼 (컨텍스트 메뉴)\n"
    "• 이름 변경\n"
    "• 아이콘 변경 (이모지 입력)\n"
    "• 새 Chapter 추가 (Notebook 메뉴)\n"
    "• 삭제\n\n"
    "+ 새 Notebook 버튼\n"
    "• 사이드바 하단에 위치\n"
    "• 클릭하면 이름 입력 후 새 Notebook 생성",
    Inches(8.4), Inches(1.2), Inches(4.7), Inches(5.9),
    size=12, color=DARK)


# ═══════════════════════════════════════════════════════
# SLIDE 5 : 컨텍스트 메뉴 (···)
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "③ 컨텍스트 메뉴 (··· 버튼)  —  이름 변경 / 아이콘 변경 / 추가 / 삭제",
         Inches(0.2), Inches(0.04), Inches(12), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "context")

def ctx(px, py, pw, ph): return ss_emu("context", px, py, pw, ph)

# 컨텍스트 메뉴 팝업 전체
l,t,w,h = ctx(48, 66, 134, 92)
add_spot(s, l, t, w, h)

# 각 메뉴 항목 라벨
ml, mw = l+w+Pt(10), Inches(1.5)
add_label(s, "이름 변경",      ml, t+Inches(0.02), mw)
add_label(s, "아이콘 변경",    ml, t+Inches(0.36), mw)
add_label(s, "새 Chapter",     ml, t+Inches(0.68), mw, bg_col=GREEN)

# ··· 버튼 위치 강조
l2,t2,w2,h2 = ctx(184, 70, 22, 22)
add_spot(s, l2, t2, w2, h2, color=ORANGE, lw=Pt(3.0))
add_label(s, "··· 버튼", l2-Inches(0.8), t2-Pt(24), Inches(1.0), bg_col=ORANGE)

# 우측 설명 카드
add_card(s, Inches(7.6), Inches(0.6), Inches(5.5), Inches(6.7))
add_text(s, "컨텍스트 메뉴 사용법",
         Inches(7.8), Inches(0.75), Inches(5.1), Inches(0.45),
         size=14, bold=True, color=DARK)
add_hline(s, Inches(7.8), Inches(1.18), Inches(5.1), MUTED)
add_text(s,
    "사용 방법\n"
    "Notebook 또는 Chapter 항목에 마우스를\n"
    "올리면 우측에 ··· 버튼이 나타납니다.\n"
    "클릭하면 컨텍스트 메뉴가 팝업됩니다.\n\n"
    "이름 변경\n"
    "→ 입력창이 열리고 새 이름을 입력합니다.\n"
    "   확인 버튼 또는 Enter 키로 저장.\n\n"
    "아이콘 변경\n"
    "→ 이모지를 직접 입력합니다.\n"
    "   예: 📚  🧠  💊  🏥  🔬  📋  📝\n\n"
    "새 Chapter 추가  ← Notebook 메뉴에만 있음\n"
    "→ 이름 입력 후 새 Chapter를 생성합니다.\n\n"
    "삭제\n"
    "→ 해당 항목과 모든 하위 데이터를 삭제합니다.\n"
    "   ⚠ 삭제 전 확인 팝업이 표시됩니다.\n"
    "   (Notebook 삭제 시 모든 Chapter·Leaf page 삭제)",
    Inches(7.8), Inches(1.28), Inches(5.1), Inches(5.85),
    size=12, color=DARK)


# ═══════════════════════════════════════════════════════
# SLIDE 6 : Chapter 뷰 & Leaf page 관리
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "④ Chapter 뷰  —  Leaf page 목록 & 추가",
         Inches(0.2), Inches(0.04), Inches(8), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "chapter")

# 위치 표시 브레드크럼
l,t,w,h = c(224, 65, 312, 22)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "현재 위치 (Notebook / Chapter)", l, t-Pt(24), Inches(2.9), bg_col=BLUE)

# Chapter 제목
l,t,w,h = c(224, 92, 158, 48)
add_spot(s, l, t, w, h)
add_label(s, "Chapter 제목", l, t+h+Pt(5), Inches(1.5))

# Chapter 설명 (있을 경우)
l,t,w,h = c(224, 140, 160, 16)
add_spot(s, l, t, w, h, color=MUTED, lw=Pt(1.5))
add_label(s, "Chapter 설명 텍스트", l, t+h+Pt(4), Inches(2.0), bg_col=MUTED)

# Leaf page 행
l,t,w,h = c(233, 192, 788, 33)
add_spot(s, l, t, w, h)
add_label(s, "Leaf page 항목 (클릭 → 상세 보기)", l, t-Pt(24), Inches(3.4))

# PDF 개수 · 날짜
l,t,w,h = c(1033, 198, 118, 20)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "PDF 개수 · 최종 수정일", l-Inches(0.3), t+h+Pt(5), Inches(2.2), bg_col=BLUE)

# + Leaf page 추가
l,t,w,h = c(233, 238, 130, 26)
add_spot(s, l, t, w, h, color=GREEN)
add_label(s, "+ Leaf page 추가", l, t+h+Pt(5), Inches(1.8), bg_col=GREEN)

# 우측 보충 설명
add_card(s, Inches(9.0), Inches(0.6), Inches(4.1), Inches(3.0))
add_text(s,
    "Leaf page 항목 행 구성\n\n"
    "• 아이콘 + 제목 (왼쪽)\n"
    "• PDF 개수 · 최종 수정일 (오른쪽)\n\n"
    "클릭하면 해당 Leaf page 상세 뷰로 이동",
    Inches(9.15), Inches(0.75), Inches(3.8), Inches(2.7),
    size=12, color=DARK)


# ═══════════════════════════════════════════════════════
# SLIDE 7 : Leaf page 뷰 — 메모 & 태그
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "⑤ Leaf page 상세 뷰  —  제목 / 태그 / 메모",
         Inches(0.2), Inches(0.04), Inches(9), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "note")

def n(px, py, pw, ph): return ss_emu("note", px, py, pw, ph)

# 브레드크럼
l,t,w,h = n(143, 58, 408, 17)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "현재 위치 (Notebook / Chapter / Leaf page)", l, t-Pt(22), Inches(4.0), bg_col=BLUE)

# 제목
l,t,w,h = n(143, 75, 315, 37)
add_spot(s, l, t, w, h)
add_label(s, "Leaf page 제목", l, t+h+Pt(5), Inches(1.8))

# 태그
l,t,w,h = n(143, 113, 180, 18)
add_spot(s, l, t, w, h, color=GREEN)
add_label(s, "태그 목록 (예: 연구, ADHD, 논문리뷰)", l+w+Pt(5), t, Inches(3.2), bg_col=GREEN)

# 메모 입력란
l,t,w,h = n(143, 132, 622, 73)
add_spot(s, l, t, w, h)
add_label(s, "메모 입력란 (자유 텍스트, 여러 줄 입력 가능)", l, t-Pt(24), Inches(4.2))

# 메모 저장
l,t,w,h = n(143, 207, 78, 24)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "메모 저장\n(클릭 시 즉시 저장)", l-Inches(0.1), t+h+Pt(5), Inches(1.6), bg_col=BLUE)

# 삭제
l2,t2,w2,h2 = n(226, 207, 55, 24)
add_spot(s, l2, t2, w2, h2)
add_label(s, "Leaf page 삭제\n(확인 후 삭제)", l2-Inches(0.1), t2+h2+Pt(5), Inches(1.6))

# 우측 설명
add_card(s, Inches(9.5), Inches(0.6), Inches(3.6), Inches(4.5))
add_text(s,
    "메모 작성 안내\n\n"
    "• 메모란은 자유 텍스트 입력\n"
    "• '메모 저장' 버튼 클릭 시 저장됨\n\n"
    "태그\n"
    "• 검색 시 태그 내용도 검색 대상에 포함됨\n\n"
    "삭제\n"
    "• Leaf page 전체를 삭제합니다\n"
    "• 삭제 전 확인 팝업 표시\n"
    "• 연결된 PDF 정보도 모두 삭제됨",
    Inches(9.65), Inches(0.75), Inches(3.3), Inches(4.2),
    size=12, color=DARK)


# ═══════════════════════════════════════════════════════
# SLIDE 8 : PDF 추가 & 뷰어
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "⑥ PDF 추가 & 뷰어",
         Inches(0.2), Inches(0.04), Inches(6), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "note")

# PDF 섹션 헤더
l,t,w,h = n(143, 232, 65, 19)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "PDF 개수", l, t-Pt(22), Inches(1.1), bg_col=BLUE)

# + PDF 추가
l,t,w,h = n(204, 231, 72, 20)
add_spot(s, l, t, w, h, color=GREEN)
add_label(s, "+ PDF 추가\n(클릭 → 파일 선택 창)", l, t+h+Pt(5), Inches(1.7), bg_col=GREEN)

# PDF 카드 헤더 (아코디언)
l,t,w,h = n(160, 260, 596, 27)
add_spot(s, l, t, w, h, color=BLUE)
add_label(s, "PDF 카드 헤더 (클릭 → 뷰어 펼치기 / 접기)", l, t-Pt(24), Inches(4.0), bg_col=BLUE)

# 코멘트 수 배지
l2,t2,w2,h2 = n(677, 263, 73, 17)
add_spot(s, l2, t2, w2, h2, color=ORANGE)
add_label(s, "코멘트 수", l2+w2+Pt(5), t2, Inches(1.1), bg_col=ORANGE)

# PDF 뷰어
l,t,w,h = n(160, 288, 398, 317)
add_spot(s, l, t, w, h)
add_label(s, "PDF 뷰어\n페이지 이동 / 줌 조절 가능", l, t+h+Pt(5), Inches(2.8))

# 우측 설명
add_card(s, Inches(9.5), Inches(0.6), Inches(3.6), Inches(5.5))
add_text(s,
    "PDF 추가 방법\n\n"
    "+ PDF 추가 버튼 클릭 →\n"
    "파일 선택 다이얼로그 열림 →\n"
    "로컬 PDF 파일 선택 →\n"
    "자동으로 목록에 추가됨\n\n"
    "PDF 카드 조작\n\n"
    "• 카드 헤더 클릭: 뷰어 펼치기/접기\n"
    "• 한 Leaf page에 PDF 여러 개 가능\n\n"
    "뷰어 조작\n\n"
    "• ◀ ▶  : 페이지 이동\n"
    "• + / − : 줌 인/아웃\n"
    "• 현재 페이지 번호 표시",
    Inches(9.65), Inches(0.75), Inches(3.3), Inches(5.2),
    size=12, color=DARK)


# ═══════════════════════════════════════════════════════
# SLIDE 9 : PDF 페이지 코멘트
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "⑦ PDF 페이지 코멘트",
         Inches(0.2), Inches(0.04), Inches(6), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "note")

# 코멘트 패널 전체
l,t,w,h = n(562, 281, 218, 325)
add_spot(s, l, t, w, h)
add_label(s, "PDF 페이지 코멘트 패널", l, t-Pt(24), Inches(2.8))

# 페이지 배지
l2,t2,w2,h2 = n(567, 298, 65, 16)
add_spot(s, l2, t2, w2, h2, color=GREEN)
add_label(s, "Page 번호\n(클릭→해당 페이지로 이동)", l2+w2+Pt(6), t2, Inches(2.4), bg_col=GREEN)

# 코멘트 본문
l3,t3,w3,h3 = n(567, 315, 196, 50)
add_spot(s, l3, t3, w3, h3, color=ACCENT)
add_label(s, "코멘트 본문", l3+w3+Pt(6), t3, Inches(1.4), bg_col=ACCENT)

# 날짜
l4,t4,w4,h4 = n(567, 368, 196, 14)
add_spot(s, l4, t4, w4, h4, color=MUTED, lw=Pt(1.5))
add_label(s, "작성 일시", l4+w4+Pt(6), t4, Inches(1.1), bg_col=MUTED)

# 코멘트 입력란
l5,t5,w5,h5 = n(562, 565, 215, 35)
add_spot(s, l5, t5, w5, h5, color=BLUE)
add_label(s, "코멘트 입력란\n(페이지 번호 + 텍스트)", l5-Inches(0.3), t5+h5+Pt(5), Inches(2.5), bg_col=BLUE)

# 하단 팁 카드
add_card(s, Inches(0.1), Inches(5.45), Inches(5.8), Inches(1.9),
         fill=RGBColor(0xFF, 0xF8, 0xEE), border=ACCENT)
add_text(s,
    "💡  코멘트 사용 팁\n\n"
    "• 현재 뷰어에 열린 페이지 번호가 입력란에 자동 반영됩니다.\n"
    "• 저장된 코멘트의 'Page N' 링크를 클릭하면 뷰어가 해당 페이지로 이동합니다.\n"
    "• 여러 PDF의 코멘트가 하나의 패널에 통합 표시됩니다.",
    Inches(0.25), Inches(5.55), Inches(5.5), Inches(1.75),
    size=11, color=DARK)


# ═══════════════════════════════════════════════════════
# SLIDE 10 : 검색 기능
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "⑧ 검색 기능  —  노트 / PDF / 태그 / 코멘트 통합 검색",
         Inches(0.2), Inches(0.04), Inches(11), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "dashboard")

# 검색창
l,t,w,h = d(1093, 2, 268, 26)
add_spot(s, l, t, w, h, color=GREEN)
add_label(s, "키워드 입력 후 Enter 키", l-Inches(1.5), t+h+Pt(6), Inches(2.5), bg_col=GREEN)

# 좌측 설명 카드
add_card(s, Inches(0.1), Inches(1.0), Inches(5.8), Inches(6.3))
add_text(s, "검색 기능 안내",
         Inches(0.25), Inches(1.1), Inches(5.5), Inches(0.45),
         size=14, bold=True, color=DARK)
add_hline(s, Inches(0.25), Inches(1.53), Inches(5.5), MUTED)
add_text(s,
    "검색 범위 (가중치 순)\n"
    "  1위  Leaf page 제목\n"
    "  2위  태그\n"
    "  3위  PDF 페이지 코멘트\n"
    "  4위  메모 본문\n"
    "  5위  PDF 파일명\n\n"
    "검색 방법\n"
    "  우측 상단 검색창에 키워드 입력 후\n"
    "  Enter 키를 누르면 검색 결과 화면으로 이동\n\n"
    "검색 결과 화면 구성\n"
    "  ─────────────────────\n"
    "  Notebooks  (제목 일치)\n"
    "  ─────────────────────\n"
    "  Chapters   (제목 일치)\n"
    "  ─────────────────────\n"
    "  Leaf pages (관련도 점수 순 정렬)\n"
    "  ─────────────────────\n\n"
    "각 항목 클릭 → 해당 위치로 바로 이동",
    Inches(0.25), Inches(1.62), Inches(5.5), Inches(5.6),
    size=12, color=DARK)


# ═══════════════════════════════════════════════════════
# SLIDE 11 : 기타 기능
# ═══════════════════════════════════════════════════════
s = blank()
add_bg(s, BEIGE)
add_text(s, "⑨ 기타 기능  —  사이드바 토글 / 언어 설정 / 데이터 저장",
         Inches(0.2), Inches(0.04), Inches(11), Inches(0.44),
         size=17, bold=True, color=DARK)
add_image(s, "chapter")

# 사이드바 접기 버튼
l,t,w,h = c(207, 392, 10, 55)
add_spot(s, l, t, Inches(0.14), h, color=BLUE)
add_label(s, "◀ 사이드바 접기", int(l-Inches(1.5)), t, Inches(1.6), bg_col=BLUE)

# 언어 선택
l2,t2,w2,h2 = c(1312, 2, 52, 26)
add_spot(s, l2, t2, w2, h2, color=GREEN)
add_label(s, "언어 설정 KO/EN", l2-Inches(0.8), t2+h2+Pt(6), Inches(1.9), bg_col=GREEN)

# 우측 설명 카드들
add_card(s, Inches(7.5), Inches(0.58), Inches(5.6), Inches(2.15))
add_text(s, "사이드바 접기 / 펼치기",
         Inches(7.65), Inches(0.68), Inches(5.3), Inches(0.38),
         size=13, bold=True, color=DARK)
add_text(s,
    "• 사이드바 우측 끝 ◀ 버튼 클릭 → 사이드바 닫힘\n"
    "• 닫힌 상태에서 화면 좌상단 ▶ 버튼 클릭 → 다시 열림\n"
    "  (▶ 버튼은 사이드바가 닫힌 상태에서만 표시됨)",
    Inches(7.65), Inches(1.07), Inches(5.3), Inches(1.55),
    size=12, color=DARK)

add_card(s, Inches(7.5), Inches(2.85), Inches(5.6), Inches(1.85))
add_text(s, "언어 설정",
         Inches(7.65), Inches(2.95), Inches(5.3), Inches(0.38),
         size=13, bold=True, color=DARK)
add_text(s,
    "• 우측 상단 KO / EN 드롭다운으로 인터페이스 언어 전환\n"
    "• 전환 즉시 메뉴·버튼 텍스트가 변경됨\n"
    "• 설정은 자동으로 저장됨",
    Inches(7.65), Inches(3.34), Inches(5.3), Inches(1.3),
    size=12, color=DARK)

add_card(s, Inches(7.5), Inches(4.8), Inches(5.6), Inches(2.55))
add_text(s, "데이터 저장 & 백업",
         Inches(7.65), Inches(4.9), Inches(5.3), Inches(0.38),
         size=13, bold=True, color=DARK)
add_text(s,
    "• 모든 데이터는 로컬 JSON 파일에 자동 저장됩니다.\n\n"
    "저장 위치\n"
    "  OrganizerData / data / app_data.json\n"
    "  (앱이 설치된 폴더 또는 시스템 사용자 데이터 폴더)\n\n"
    "자동 백업\n"
    "  데이터 저장마다 타임스탬프 백업 파일 자동 생성\n"
    "  경로: OrganizerData / data / backups /",
    Inches(7.65), Inches(5.29), Inches(5.3), Inches(2.0),
    size=12, color=DARK)


# ═══════════════════════════════════════════════════════
# PPTX 저장
# ═══════════════════════════════════════════════════════
OUT = "/mnt/c/EunYoungKim/Organizer/(ref)manual.pptx"
prs.save(OUT)
print("Saved PPTX:", OUT)
print(f"Total slides: {len(prs.slides)}")
